from http.server import BaseHTTPRequestHandler
import json, io, base64
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY_DARK = RGBColor(0x09, 0x21, 0x41)
NAVY_MED  = RGBColor(0x0B, 0x31, 0x5F)
ORANGE    = RGBColor(0xEC, 0x69, 0x07)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0xCC, 0xD5, 0xE8)

def set_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_bar(slide, W, color=None, height=160000):
    bar = slide.shapes.add_shape(1, 0, 0, W, Emu(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color or ORANGE
    bar.line.fill.background()

def add_tb(slide, text, x, y, w, h, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = 'Century Gothic'
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or WHITE
    return tb

def add_separator(slide, y, W):
    ln = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Emu(25000))
    ln.fill.solid(); ln.fill.fore_color.rgb = ORANGE; ln.line.fill.background()

def slide_text(prs, title, paragraphs, subtitle=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    set_bg(sl, NAVY_DARK)
    add_bar(sl, W)
    add_tb(sl, title.upper(), Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.75), size=22, bold=True)
    if subtitle:
        add_tb(sl, subtitle, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.45), size=12, color=GRAY_TEXT)
    add_separator(sl, Inches(1.1) if not subtitle else Inches(1.4), W)
    y_start = Inches(1.25) if not subtitle else Inches(1.55)
    tb = sl.shapes.add_textbox(Inches(0.5), y_start, Inches(12.3), Inches(5.7))
    tf = tb.text_frame; tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para.get('text', '')
        p.font.name = 'Century Gothic'
        p.font.size = Pt(para.get('size', 13))
        p.font.bold = para.get('bold', False)
        p.font.color.rgb = WHITE if para.get('bold') else GRAY_TEXT
        p.space_after = Pt(para.get('space', 5))
    return sl

def slide_two_col(prs, title, left_title, left_items, right_title, right_items):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    set_bg(sl, NAVY_DARK)
    add_bar(sl, W)
    add_tb(sl, title.upper(), Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.75), size=22, bold=True)
    add_separator(sl, Inches(1.1), W)
    # Left
    add_tb(sl, left_title.upper(), Inches(0.5), Inches(1.2), Inches(5.9), Inches(0.45), size=11, bold=True, color=ORANGE)
    tb_l = sl.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(5.9), Inches(5.3))
    tf_l = tb_l.text_frame; tf_l.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = '• ' + item; p.font.name = 'Century Gothic'; p.font.size = Pt(12); p.font.color.rgb = GRAY_TEXT; p.space_after = Pt(7)
    # Divider
    div = sl.shapes.add_shape(1, Inches(6.55), Inches(1.2), Emu(25000), Inches(5.5))
    div.fill.solid(); div.fill.fore_color.rgb = NAVY_MED; div.line.fill.background()
    # Right
    add_tb(sl, right_title.upper(), Inches(6.7), Inches(1.2), Inches(5.9), Inches(0.45), size=11, bold=True, color=ORANGE)
    tb_r = sl.shapes.add_textbox(Inches(6.7), Inches(1.7), Inches(5.9), Inches(5.3))
    tf_r = tb_r.text_frame; tf_r.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.text = '• ' + item; p.font.name = 'Century Gothic'; p.font.size = Pt(12); p.font.color.rgb = GRAY_TEXT; p.space_after = Pt(7)
    return sl

def slide_price_table(prs, groups, currency='EUR', total=0):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    set_bg(sl, NAVY_DARK)
    add_bar(sl, W)
    add_tb(sl, 'PROPUNERE COMERCIALĂ', Inches(0.5), Inches(0.22), Inches(10), Inches(0.75), size=22, bold=True)
    add_tb(sl, currency, Inches(11.5), Inches(0.22), Inches(1.2), Inches(0.75), size=13, color=ORANGE, align=PP_ALIGN.RIGHT)
    y = Inches(1.1)
    for group in groups:
        if not group.get('rows'): continue
        headers = group['headers']
        rows = group['rows']
        label = group['label']
        n_cols = len(headers)
        n_rows = min(len(rows) + 1, 15)
        # Group label bar
        lbar = sl.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Emu(270000))
        lbar.fill.solid(); lbar.fill.fore_color.rgb = NAVY_MED; lbar.line.fill.background()
        add_tb(sl, label.upper(), Inches(0.6), y + Emu(55000), Inches(6), Emu(180000), size=10, bold=True, color=ORANGE)
        y += Emu(300000)
        tbl_h = Emu(n_rows * 310000)
        tbl = sl.shapes.add_table(n_rows, n_cols, Inches(0.5), y, Inches(12.3), tbl_h).table
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = ORANGE
            p = cell.text_frame.paragraphs[0]
            p.text = h; p.font.name = 'Century Gothic'; p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = WHITE
        for ri, row in enumerate(rows[:n_rows-1]):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri+1, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY_MED if ri % 2 == 0 else NAVY_DARK
                p = cell.text_frame.paragraphs[0]
                p.text = str(val); p.font.name = 'Century Gothic'; p.font.size = Pt(9); p.font.color.rgb = WHITE
        y += tbl_h + Emu(180000)
    # Total
    if total:
        tot_bar = sl.shapes.add_shape(1, Inches(7.5), y, Inches(5.3), Emu(350000))
        tot_bar.fill.solid(); tot_bar.fill.fore_color.rgb = ORANGE; tot_bar.line.fill.background()
        add_tb(sl, 'TOTAL: {:,.2f} {}'.format(total, currency), Inches(7.6), y + Emu(70000), Inches(5), Emu(230000), size=14, bold=True, color=WHITE)
    return sl

def slide_terms(prs, validity, payment, terms):
    paras = []
    if validity: paras.append({'text': 'Valabilitate ofertă: ' + validity, 'bold': True, 'size': 14, 'space': 10})
    if payment: paras.append({'text': 'Termeni de plată: ' + payment, 'bold': True, 'size': 14, 'space': 10})
    if terms:
        paras.append({'text': '', 'size': 5, 'space': 3})
        paras.append({'text': 'Precizări și Assumptions:', 'bold': True, 'size': 13, 'space': 8})
        for t in terms:
            if t.strip(): paras.append({'text': '• ' + t.strip(), 'size': 12, 'space': 5})
    return slide_text(prs, 'Termeni Comerciali & Assumptions', paras)

def generate_offer_pptx(data):
    template_b64 = data.get('template')
    keep_slides  = data.get('keepSlides', [])
    offer        = data.get('offer', {})
    pl           = data.get('pl', {})
    settings     = data.get('settings', {})

    # Load template & keep selected slides
    prs = Presentation(io.BytesIO(base64.b64decode(template_b64)))
    sldIdLst = prs.slides._sldIdLst
    all_ids  = list(sldIdLst)
    keep_set = set(keep_slides)
    for i, sldId in enumerate(all_ids):
        if i not in keep_set:
            sldIdLst.remove(sldId)

    currency  = pl.get('currency', 'EUR')
    pl_rows   = pl.get('rows', [])

    # 1. Intelegerea cerintelor
    needs = offer.get('clientNeeds', '')
    if needs:
        paras = [{'text': needs, 'size': 14, 'space': 12}]
        for ch in offer.get('challenges', []):
            if ch.strip(): paras.append({'text': '• ' + ch, 'size': 13, 'space': 6})
        slide_text(prs, 'Înțelegerea Cerințelor', paras, subtitle=offer.get('clientName',''))

    # 2. Solutia propusa / abordare
    approach = offer.get('approach', '')
    solution = offer.get('solution', '')
    if approach or solution:
        left  = [x.strip() for x in approach.split('\n') if x.strip()]
        right = [x.strip() for x in solution.split('\n') if x.strip()]
        if left or right:
            slide_two_col(prs, 'Soluția Propusă', 'Abordare & Metodologie', left[:7], 'Beneficii Cheie', right[:7])

    # 3. Plan implementare
    timeline = offer.get('timeline', '')
    if timeline:
        paras = [{'text': t.strip(), 'size': 13, 'space': 7} for t in timeline.split('\n') if t.strip()]
        slide_text(prs, 'Plan de Implementare', paras)

    # 4. Echipa
    team = offer.get('team', '')
    if team:
        paras = [{'text': t.strip(), 'size': 13, 'space': 7} for t in team.split('\n') if t.strip()]
        slide_text(prs, 'Echipa Dedicată', paras)

    # 5. Tabel preturi
    if pl_rows:
        type_cfg = {
            'hw':  {'label': 'Echipamente Hardware',      'headers': ['Denumire', 'Cod produs', 'Cant.', 'Preț unitar', 'Total']},
            'sw':  {'label': 'Licențe & Software',        'headers': ['Produs', 'SKU', 'Cant.', 'Preț unitar', 'Total']},
            'svc': {'label': 'Servicii Profesionale',     'headers': ['Serviciu / Rol', 'Zile', 'Tarif/zi', 'Total', '']},
            'cl':  {'label': 'Cloud Services',            'headers': ['Serviciu', 'Vendor', 'Cost/lună', 'Luni', 'Total estimat']},
        }
        groups_map = {}
        total_sell = 0
        for row in pl_rows:
            t = row.get('type', 'hw')
            if t not in groups_map: groups_map[t] = []
            sell = float(row.get('sellPrice', 0))
            qty  = float(row.get('qty', 1))
            total_sell += sell
            unit_sell = sell / qty if qty else sell
            groups_map[t].append([
                row.get('name','')[:45],
                row.get('code',''),
                int(qty),
                '{:,.2f}'.format(unit_sell),
                '{:,.2f}'.format(sell)
            ])
        price_groups = []
        for t, cfg in type_cfg.items():
            if t in groups_map:
                price_groups.append({'type': t, 'label': cfg['label'], 'headers': cfg['headers'], 'rows': groups_map[t]})
        if price_groups:
            slide_price_table(prs, price_groups, currency, total_sell)

    # 6. Termeni
    validity = settings.get('validity', '30 de zile calendaristice')
    payment  = settings.get('payment', '')
    terms    = []
    if settings.get('termsDelivery'):    terms.append(settings['termsDelivery'])
    if settings.get('termsWarranty'):    terms.append(settings['termsWarranty'])
    if settings.get('termsObservations'):terms.append(settings['termsObservations'])
    for a in offer.get('assumptions', []):
        if a.strip(): terms.append(a)
    slide_terms(prs, validity, payment, terms)

    buf = io.BytesIO()
    prs.save(buf)
    return base64.b64encode(buf.getvalue()).decode('utf-8')

class handler(BaseHTTPRequestHandler):
    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def do_POST(self):
        length = int(self.headers['Content-Length'])
        data   = json.loads(self.rfile.read(length))
        try:
            pptx_b64 = generate_offer_pptx(data)
            self.send_response(200)
            self.send_header('Content-Type', 'application/json')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()
            self.wfile.write(json.dumps({'pptx': pptx_b64}).encode())
        except Exception as e:
            import traceback
            self.send_response(500)
            self.send_header('Content-Type', 'application/json')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()
            self.wfile.write(json.dumps({'error': str(e), 'trace': traceback.format_exc()}).encode())
